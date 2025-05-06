"""
Generates a PowerPoint presentation summarizing the eco-travel analysis.
"""

import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
import pandas as pd

# Define the output directory for the presentation
PPTX_DIR = "results"
PPTX_FILENAME = "eco_travel_analysis_presentation.pptx"

def setup_presentation_directory():
    """Creates the presentation output directory if it doesn't exist."""
    if not os.path.exists(PPTX_DIR):
        os.makedirs(PPTX_DIR)
    print(f"Presentation directory '{PPTX_DIR}' is ready.")

def add_title_slide(prs, title, subtitle):
    """Adds a title slide to the presentation.

    Args:
        prs (Presentation): The presentation object.
        title (str): The main title for the slide.
        subtitle (str): The subtitle for the slide.
    """
    title_slide_layout = prs.slide_layouts[0] # Title Slide layout
    slide = prs.slides.add_slide(title_slide_layout)
    title_placeholder = slide.shapes.title
    subtitle_placeholder = slide.placeholders[1]

    title_placeholder.text = title
    subtitle_placeholder.text = subtitle

def add_content_slide(prs, title, content_list):
    """Adds a content slide (Title and Content layout) to the presentation.

    Args:
        prs (Presentation): The presentation object.
        title (str): The title for the slide.
        content_list (list): A list of strings, each representing a bullet point.
    """
    content_slide_layout = prs.slide_layouts[1] # Title and Content layout
    slide = prs.slides.add_slide(content_slide_layout)
    title_placeholder = slide.shapes.title
    body_placeholder = slide.placeholders[1]

    title_placeholder.text = title
    tf = body_placeholder.text_frame
    tf.clear() # Clear existing text

    for item in content_list:
        p = tf.add_paragraph()
        p.text = item
        p.level = 0 # Top level bullet
        p.font.size = Pt(18)

def add_image_slide(prs, title, image_path, notes=""):
    """Adds a slide with a title and an image.

    Args:
        prs (Presentation): The presentation object.
        title (str): The title for the slide.
        image_path (str): Path to the image file.
        notes (str, optional): Speaker notes for the slide. Defaults to "".
    """
    blank_slide_layout = prs.slide_layouts[6] # Blank layout
    slide = prs.slides.add_slide(blank_slide_layout)

    # Add title
    title_shape = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(9), Inches(0.8))
    tf = title_shape.text_frame
    p = tf.add_paragraph()
    p.text = title
    p.font.size = Pt(32)
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER

    # Add image, centered
    if os.path.exists(image_path):
        img_width_in = 8 # Desired image width in inches
        img_height_in = 5 # Desired image height in inches
        left = (prs.slide_width.inches - img_width_in) / 2
        top = Inches(1.2) # Position below title
        slide.shapes.add_picture(image_path, left, top, width=Inches(img_width_in), height=Inches(img_height_in))
    else:
        print(f"Warning: Image not found at {image_path}")
        # Add a placeholder text if image is missing
        txt_box = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(1))
        tf = txt_box.text_frame
        p = tf.add_paragraph()
        p.text = f"Image not found: {os.path.basename(image_path)}"
        p.font.color.rgb = RGBColor(255, 0, 0) # Red text
        p.alignment = PP_ALIGN.CENTER

    # Add speaker notes if provided
    if notes:
        notes_slide = slide.notes_slide
        text_frame = notes_slide.notes_text_frame
        text_frame.text = notes

def generate_presentation(df_results, key_findings, routes):
    """Generates the PowerPoint presentation from analysis results.

    Args:
        df_results (pd.DataFrame): DataFrame containing the scenario analysis results.
        key_findings (list): List of strings summarizing key findings.
        routes (dict): Dictionary containing route details.
    """
    setup_presentation_directory()
    prs = Presentation()
    prs.slide_width = Inches(16)
    prs.slide_height = Inches(9)

    # --- Slide 1: Title ---
    add_title_slide(prs,
                    "Eco-Friendly Travel: Grenoble to Abuja",
                    "An Environmental Analysis of Travel Options")

    # --- Slide 2: Introduction ---
    intro_content = [
        "Departure: Lycée Pierre Termier – Grenoble – France",
        "Arrival: Abuja, Nigeria",
        "Distance: ~4,019 km (direct)",
        "Objective: Analyze cost, duration, and carbon footprint of travel options.",
        "Scenarios: 1-week vs. 1-month vacation (including travel time)"
    ]
    add_content_slide(prs, "Introduction", intro_content)

    # --- Slide 3: Transportation Options Summary ---
    options_content = []
    for name, data in routes.items():
        options_content.append(f"{name}: ~{data['travel_time_hours']:.0f} hrs one-way, {data['total_carbon_one_way']:.1f} kg CO2e one-way, €{data['cost_eur']/2:.0f} one-way")
    add_content_slide(prs, "Transportation Options Overview", options_content)

    # --- Slide 4: Carbon Footprint Comparison Plot ---
    carbon_plot_path = os.path.join("results", "visualizations", "carbon_footprint_comparison.png")
    add_image_slide(prs, "Carbon Footprint Comparison (Round Trip)", carbon_plot_path,
                    notes="This chart compares the total CO2 equivalent emissions for each route under both scenarios. Note the significant difference between air travel and the land/sea options.")

    # --- Slide 5: Carbon Breakdown Plot ---
    breakdown_plot_path = os.path.join("results", "visualizations", "carbon_breakdown.png")
    add_image_slide(prs, "Carbon Footprint Breakdown (One-Way)", breakdown_plot_path,
                    notes="This shows which modes of transport contribute most to the carbon footprint for each route. For air travel, the flight itself dominates. For others, contributions are more varied.")

    # --- Slide 6: Time Efficiency & Feasibility ---
    time_plot_path = os.path.join("results", "visualizations", "time_distribution.png")
    add_image_slide(prs, "Time Distribution: Travel vs. Destination", time_plot_path,
                    notes="This visualizes how total vacation time is split between traveling and being at the destination. Overland routes consume much more travel time, making them infeasible for the 1-week scenario.")

    # --- Slide 7: Carbon per Vacation Day Plot ---
    carbon_per_day_plot_path = os.path.join("results", "visualizations", "carbon_per_vacation_day.png")
    add_image_slide(prs, "Carbon Footprint per Vacation Day", carbon_per_day_plot_path,
                    notes="This metric shows the carbon 'cost' per day actually spent enjoying the destination. Longer trips with low-carbon transport become much more efficient by this measure.")

    # --- Slide 8: Cost Comparison (Hostel) Plot ---
    cost_hostel_plot_path = os.path.join("results", "visualizations", "cost_comparison_hostel.png")
    add_image_slide(prs, "Total Cost Comparison (Hostel)", cost_hostel_plot_path,
                    notes="Compares the estimated total trip cost (transport + hostel accommodation) for feasible options. Costs are relatively similar for the 1-month scenario across different routes.")

    # --- Slide 9: Key Findings / Summary ---
    add_content_slide(prs, "Key Findings", key_findings)

    # --- Slide 10: Recommendations ---
    reco_content = []
    # Extract recommendations from findings (assuming they are the last items)
    reco_start_index = -1
    for i, finding in enumerate(key_findings):
        if finding.startswith("5. Environmental Recommendation Summary:"):
            reco_start_index = i
            break
    if reco_start_index != -1:
        reco_content = [f.replace("   - ", "") for f in key_findings[reco_start_index+1:]]
    else:
        reco_content = ["Recommendations could not be automatically extracted."]

    add_content_slide(prs, "Recommendations", reco_content)

    # --- Slide 11: Conclusion ---
    # Extract conclusion from markdown or generate dynamically if needed
    # For now, using a placeholder based on findings
    conclusion_content = [
        "1-Week Trip: Flying is the only feasible option. Focus on mitigation.",
        "1-Month Trip: Land/Sea routes offer >80% carbon reduction vs. flying.",
        "'Land & Sea' is the most carbon-efficient overall.",
        "Trade-offs: Time, cost, complexity vs. environmental impact.",
        "Longer travel time enables significantly greener choices."
    ]
    add_content_slide(prs, "Conclusion", conclusion_content)

    # --- Slide 12: Q&A ---
    add_title_slide(prs, "Q&A", "Thank you!")

    # --- Save Presentation ---
    pptx_path = os.path.join(PPTX_DIR, PPTX_FILENAME)
    try:
        prs.save(pptx_path)
        print(f"Presentation successfully saved to: {pptx_path}")
    except Exception as e:
        print(f"Error saving presentation: {e}")

# Example usage (if run directly, though typically called from main script)
if __name__ == '__main__':
    # Create dummy data for testing if needed
    print("This script is intended to be imported and called from main.py")
    # Example: generate_presentation(dummy_df, dummy_findings, dummy_routes)